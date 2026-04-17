#!/usr/bin/env python3
"""
Build the GigShield Phase 3 pitch deck as a PPTX.
Guidewire UI style — navy + orange, 16:9, 18 slides, mirrors gigshield-pitch.html.

Run:  python3 build_pptx.py
Out:  gigshield-pitch.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# --- Guidewire palette ----------------------------------------------------
GW_NAVY    = RGBColor(0x00, 0x2B, 0x5C)
GW_NAVY_2  = RGBColor(0x01, 0x32, 0x62)
GW_ORANGE  = RGBColor(0xF2, 0x65, 0x22)
GW_ORANGE2 = RGBColor(0xFF, 0x81, 0x39)
GW_BLUE    = RGBColor(0x00, 0xA7, 0xE1)
GW_GREEN   = RGBColor(0x1F, 0xAE, 0x5F)
GW_RED     = RGBColor(0xD1, 0x34, 0x38)
GW_AMBER   = RGBColor(0xE8, 0xA3, 0x3D)
GW_GRAY1   = RGBColor(0x5A, 0x67, 0x74)
GW_GRAY2   = RGBColor(0x8A, 0x94, 0xA0)
GW_BG      = RGBColor(0xF5, 0xF7, 0xFA)
GW_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GW_INK     = RGBColor(0x0E, 0x1A, 0x2A)

# 16:9 slide size: 13.333 x 7.5 inches
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Logo assets — copied from packages/frontend/public/logos. Two variants:
#   - gigshield.png      : navy shield + navy wordmark, for light backgrounds
#   - gigshield-dark.png : white/pale shield + wordmark, for dark (navy) backgrounds
ASSET_DIR       = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'assets')
LOGO_LIGHT_BG   = os.path.join(ASSET_DIR, 'gigshield.png')       # use on white
LOGO_DARK_BG    = os.path.join(ASSET_DIR, 'gigshield-dark.png')  # use on navy

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]

# --- Helpers --------------------------------------------------------------
def add_logo(slide, left, top, *, height, on_dark=False):
    """Insert the GigShield logo PNG, preserving aspect ratio.

    `height` is the desired rendered height (EMU / Inches object).
    Pass `on_dark=True` to use the white-on-dark variant (for navy backgrounds).
    """
    path = LOGO_DARK_BG if on_dark else LOGO_LIGHT_BG
    # python-pptx will infer width from the PNG's aspect ratio when width is omitted.
    return slide.shapes.add_picture(path, left, top, height=height)

def add_rect(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp

def add_text(slide, left, top, width, height, text, *,
             font='Inter', size=18, bold=False, color=GW_INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def add_runs(slide, left, top, width, height, runs, *,
             font='Inter', align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15):
    """runs = list of (text, {size, bold, color, italic})"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for text, attrs in runs:
        if '\n' in text:
            parts = text.split('\n')
            for i, part in enumerate(parts):
                if i > 0:
                    p = tf.add_paragraph()
                    p.alignment = align
                    p.line_spacing = line_spacing
                r = p.add_run()
                r.text = part
                r.font.name = attrs.get('font', font)
                r.font.size = Pt(attrs.get('size', 14))
                r.font.bold = attrs.get('bold', False)
                r.font.italic = attrs.get('italic', False)
                r.font.color.rgb = attrs.get('color', GW_INK)
        else:
            r = p.add_run()
            r.text = text
            r.font.name = attrs.get('font', font)
            r.font.size = Pt(attrs.get('size', 14))
            r.font.bold = attrs.get('bold', False)
            r.font.italic = attrs.get('italic', False)
            r.font.color.rgb = attrs.get('color', GW_INK)
    return tb

def add_slide_chrome(slide, page_num, section_name):
    """Navy header bar + footer that every content slide has."""
    # Background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, GW_WHITE)
    # Header bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.55), GW_NAVY)
    # GigShield logo (white variant) in header — the PNG is a shield + wordmark
    # stacked, so we inset slightly and keep it small so it reads as a corner
    # mark, not the main title.
    add_logo(slide, Inches(0.35), Inches(0.06), height=Inches(0.45), on_dark=True)
    # Orange accent divider right after the logo
    add_rect(slide, Inches(1.85), Inches(0.18), Inches(0.02), Inches(0.2), GW_ORANGE)
    # Brand text
    add_text(slide, Inches(1.95), Inches(0.13), Inches(8), Inches(0.3),
             f"Phase 3 Final Submission   ·   Guidewire DEVTrails 2026",
             size=11, bold=True, color=RGBColor(0xD8, 0xE2, 0xF0))
    # Section on the right
    add_text(slide, Inches(9.5), Inches(0.13), Inches(3.6), Inches(0.3),
             section_name.upper(),
             size=10, bold=True, color=RGBColor(0xC8, 0xD3, 0xE3), align=PP_ALIGN.RIGHT)
    # Footer
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.02), RGBColor(0xDC, 0xE2, 0xEA))
    add_text(slide, Inches(0.45), Inches(7.22), Inches(8), Inches(0.25),
             "GIGSHIELD © 2026 TEAM VOID", size=9, bold=True, color=GW_GRAY2)
    # Page num pill
    pg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(12.45), Inches(7.2), Inches(0.55), Inches(0.25))
    pg.fill.solid(); pg.fill.fore_color.rgb = GW_NAVY; pg.line.fill.background()
    pg.adjustments[0] = 0.3
    add_text(slide, Inches(12.45), Inches(7.24), Inches(0.55), Inches(0.25),
             f"{page_num:02d}", size=9, bold=True, color=GW_WHITE, align=PP_ALIGN.CENTER)

def add_eyebrow(slide, top, text):
    # Small orange bar + eyebrow text
    bar = add_rect(slide, Inches(0.7), top + Inches(0.09), Inches(0.25), Inches(0.03), GW_ORANGE)
    add_text(slide, Inches(1.05), top, Inches(8), Inches(0.25),
             text.upper(), size=10, bold=True, color=GW_ORANGE)

def add_title(slide, top, title_text, accent_word=None):
    """Big slide title with optional orange accent word."""
    if accent_word and accent_word in title_text:
        parts = title_text.split(accent_word)
        runs = []
        runs.append((parts[0], {'size': 32, 'bold': True, 'color': GW_NAVY}))
        runs.append((accent_word, {'size': 32, 'bold': True, 'color': GW_ORANGE}))
        if len(parts) > 1:
            runs.append((parts[1], {'size': 32, 'bold': True, 'color': GW_NAVY}))
        add_runs(slide, Inches(0.7), top, Inches(12), Inches(1.3), runs)
    else:
        add_text(slide, Inches(0.7), top, Inches(12), Inches(1.3),
                 title_text, size=32, bold=True, color=GW_NAVY, line_spacing=1.08)

def add_subtitle(slide, top, text):
    add_text(slide, Inches(0.7), top, Inches(11), Inches(0.7),
             text, size=14, color=GW_GRAY1, line_spacing=1.35)

def add_kpi_card(slide, left, top, width, height, *,
                 kpi, unit, label, desc=None,
                 accent=GW_ORANGE, navy=False):
    """Guidewire-style KPI card with orange top-border."""
    bg = GW_NAVY if navy else GW_WHITE
    # Border (top bar)
    add_rect(slide, left, top, width, Inches(0.07), accent)
    # Body
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.07),
                                  width, height - Inches(0.07))
    body.fill.solid(); body.fill.fore_color.rgb = bg
    body.line.color.rgb = RGBColor(0xE0, 0xE6, 0xEE) if not navy else GW_NAVY
    body.line.width = Pt(0.5)
    body.shadow.inherit = False
    # KPI
    kpi_color = GW_WHITE if navy else GW_NAVY
    txt_color = GW_WHITE if navy else GW_INK
    desc_color = RGBColor(0xD0, 0xD8, 0xE4) if navy else GW_GRAY1
    add_runs(slide, left + Inches(0.15), top + Inches(0.25),
             width - Inches(0.2), Inches(0.7),
             [(kpi, {'size': 28, 'bold': True, 'color': kpi_color}),
              (' ' + unit, {'size': 14, 'bold': True, 'color': accent})])
    add_text(slide, left + Inches(0.15), top + Inches(0.95),
             width - Inches(0.2), Inches(0.3),
             label, size=12, bold=True, color=txt_color)
    if desc:
        add_text(slide, left + Inches(0.15), top + Inches(1.25),
                 width - Inches(0.2), height - Inches(1.3),
                 desc, size=10, color=desc_color, line_spacing=1.3)

def add_pill(slide, left, top, text, bg_color, text_color, width=None):
    w = width or Inches(1.3)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.28))
    pill.fill.solid(); pill.fill.fore_color.rgb = bg_color; pill.line.fill.background()
    pill.adjustments[0] = 0.5
    add_text(slide, left, top + Inches(0.035), w, Inches(0.25),
             text, size=9, bold=True, color=text_color, align=PP_ALIGN.CENTER)

def add_bullets(slide, left, top, width, height, items, *, size=12, color=GW_INK):
    """Custom bullet list with orange diamond markers."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.35
        if i > 0:
            p.space_before = Pt(4)
        # bullet marker
        r0 = p.add_run(); r0.text = '◆ '
        r0.font.name = 'Inter'; r0.font.size = Pt(size - 2)
        r0.font.bold = True; r0.font.color.rgb = GW_ORANGE
        # text
        r1 = p.add_run(); r1.text = item
        r1.font.name = 'Inter'; r1.font.size = Pt(size)
        r1.font.color.rgb = color

# =========================================================================
# SLIDE 1 — TITLE
# =========================================================================
s = prs.slides.add_slide(BLANK)
# Navy gradient background (fake with two rects + radial accent)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, GW_NAVY)
# Orange radial glow top-right (simulated with a soft rect)
glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-2), Inches(6), Inches(6))
glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x00, 0x3E, 0x80)
glow.line.fill.background()
# Orange block corner
add_rect(s, Inches(11.5), 0, Inches(1.8), Inches(0.4), GW_ORANGE)

# Logo block — real GigShield logo (white-on-dark variant) top-left
add_logo(s, Inches(0.8), Inches(0.7), height=Inches(1.0), on_dark=True)
add_text(s, Inches(2.15), Inches(1.05), Inches(10), Inches(0.3),
         "GUIDEWIRE DEVTRAILS 2026   ·   PHASE 3 FINAL",
         size=10, bold=True, color=RGBColor(0xB0, 0xC0, 0xD8))

# Huge title — "GigShield."
add_runs(s, Inches(0.8), Inches(2.1), Inches(12), Inches(1.8),
         [('Gig',    {'size': 72, 'bold': True, 'color': GW_WHITE}),
          ('Shield', {'size': 72, 'bold': True, 'color': GW_ORANGE}),
          ('.',      {'size': 72, 'bold': True, 'color': GW_WHITE})])

# Tagline
add_text(s, Inches(0.8), Inches(3.4), Inches(11), Inches(1.2),
         "AI-powered parametric income protection for India's 7.7 million gig workers — "
         "zero paperwork, sub-10-minute payout, real Razorpay integration.",
         size=18, color=RGBColor(0xE0, 0xE8, 0xF2), line_spacing=1.35)

# Chip row
chips = [
    ('Parametric Insurance', GW_ORANGE, GW_WHITE),
    ('BAS Anti-Spoofing',    RGBColor(0x1A, 0x44, 0x7A), GW_WHITE),
    ('OpenAI GPT-5.4',       RGBColor(0x1A, 0x44, 0x7A), GW_WHITE),
    ('Live Razorpay API',    RGBColor(0x1A, 0x44, 0x7A), GW_WHITE),
    ('11 Indian Languages',  RGBColor(0x1A, 0x44, 0x7A), GW_WHITE),
    ('Server-Sent Events',   RGBColor(0x1A, 0x44, 0x7A), GW_WHITE),
]
x = Inches(0.8)
for label, bg, fg in chips:
    w = Inches(0.15 + len(label) * 0.085)
    add_pill(s, x, Inches(5.0), label, bg, fg, width=w)
    x = x + w + Inches(0.1)

# Meta block footer
add_rect(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.02), RGBColor(0x44, 0x66, 0x88))
add_runs(s, Inches(0.8), Inches(6.3), Inches(6), Inches(1.0),
         [('TEAM',    {'size': 9,  'bold': True, 'color': RGBColor(0x8A, 0xA0, 0xC0)}),
          ('   Team Void\n',     {'size': 12, 'bold': True, 'color': GW_WHITE}),
          ('LIVE',    {'size': 9,  'bold': True, 'color': RGBColor(0x8A, 0xA0, 0xC0)}),
          ('   gigshield.in\n',     {'size': 12, 'bold': True, 'color': GW_WHITE}),
          ('VIDEO',   {'size': 9,  'bold': True, 'color': RGBColor(0x8A, 0xA0, 0xC0)}),
          ('   youtu.be/bx8AVAU_amk',  {'size': 12, 'bold': True, 'color': GW_WHITE})])
add_runs(s, Inches(7.0), Inches(6.3), Inches(5.7), Inches(1.0),
         [('SUBMISSION', {'size': 9,  'bold': True, 'color': RGBColor(0x8A, 0xA0, 0xC0)}),
          ('   April 18, 2026\n',   {'size': 12, 'bold': True, 'color': GW_WHITE}),
          ('DEPLOYED',   {'size': 9,  'bold': True, 'color': RGBColor(0x8A, 0xA0, 0xC0)}),
          ('   Hostinger + Cloudflare\n', {'size': 12, 'bold': True, 'color': GW_WHITE}),
          ('REPO',       {'size': 9,  'bold': True, 'color': RGBColor(0x8A, 0xA0, 0xC0)}),
          ('   uroy80/TeamVoid-Guidewire-Devtrails', {'size': 12, 'bold': True, 'color': GW_WHITE})],
         align=PP_ALIGN.LEFT)

# =========================================================================
# SLIDE 2 — THE PROBLEM
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_slide_chrome(s, 2, "01 / The Problem")
add_eyebrow(s, Inches(0.85), "Problem Statement")
add_title(s, Inches(1.15), "When the weather breaks, income disappears.", accent_word="income disappears.")
add_subtitle(s, Inches(2.2),
    "Q-Commerce delivery partners earn per-order. Traditional insurance has never reached them — "
    "no payslips, no ITRs, 42–75-day claim cycles.")

# Quote block
quote_bg = add_rect(s, Inches(0.7), Inches(3.0), Inches(11.9), Inches(1.1),
                    RGBColor(0xFE, 0xF3, 0xEC))
add_rect(s, Inches(0.7), Inches(3.0), Inches(0.08), Inches(1.1), GW_ORANGE)
add_text(s, Inches(1.0), Inches(3.15), Inches(11.4), Inches(0.6),
         '"Kal 6 ghante kaam nahi kar paya waterlogging ki wajah se. ₹400 ka loss. '
         'Hamein kuch nahi milta."',
         size=16, italic=True, bold=True, color=GW_NAVY, line_spacing=1.35)
add_text(s, Inches(1.0), Inches(3.7), Inches(11.4), Inches(0.3),
         "— Delivery partner, Andheri East, Mumbai · July 2025",
         size=11, color=GW_GRAY1)

# 4 KPI cards
card_y = Inches(4.4)
card_w = Inches(2.85); gap = Inches(0.15)
card_h = Inches(2.4)

kpi_data = [
    ("3–6", "hrs",  "Mumbai monsoon",          "Waterlogging shuts down entire zones mid-shift."),
    ("₹640", "lost", "Per event, median rider", "Six hours of high-peak earning vanishes."),
    ("72",  "%",    "Income earned per-order", "No hourly floor — no deliveries = zero income."),
    ("0",   "%",    "Of claims paid today",    "No existing insurer underwrites this risk at this price."),
]
for i, (kpi, unit, label, desc) in enumerate(kpi_data):
    add_kpi_card(s, Inches(0.7) + i * (card_w + gap), card_y, card_w, card_h,
                 kpi=kpi, unit=unit, label=label, desc=desc)

# =========================================================================
# SLIDE 3 — MARKET
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_slide_chrome(s, 3, "02 / Market")
add_eyebrow(s, Inches(0.85), "Total Addressable Market")
add_title(s, Inches(1.15), "A ₹1,386 Cr annual premium pool, untapped.",
          accent_word="untapped.")
add_subtitle(s, Inches(2.2),
    "Parametric micro-insurance is the only math that clears at ₹15/month. "
    "Remove the human adjuster — unlock the entire segment.")

# Stat row — 4 stats in a single wide card
stat_y = Inches(3.0)
stat_w = Inches(2.95); stat_h = Inches(1.6)
stats = [
    ("7.7", "M",   "Indian Gig Workers (TAM)"),
    ("1.2", "M",   "Q-Commerce riders (SAM)"),
    ("1,386", "Cr",  "Annual premium pool @ ₹15/mo"),
    ("150", "K",   "Year-1 target (SOM)"),
]
for i, (big, unit, lbl) in enumerate(stats):
    left = Inches(0.7) + i * (stat_w + Inches(0.05))
    add_rect(s, left, stat_y, stat_w, stat_h, GW_WHITE, line=RGBColor(0xD0, 0xD8, 0xE4))
    add_runs(s, left + Inches(0.2), stat_y + Inches(0.25), stat_w - Inches(0.3), Inches(0.8),
             [(big, {'size': 34, 'bold': True, 'color': GW_NAVY}),
              (' ' + unit, {'size': 15, 'bold': True, 'color': GW_ORANGE})])
    add_text(s, left + Inches(0.2), stat_y + Inches(1.05), stat_w - Inches(0.3), Inches(0.5),
             lbl, size=11, bold=True, color=GW_GRAY1)

# Why now / Why incumbents can't — two columns
col_y = Inches(4.8); col_h = Inches(2.2)
col_w = Inches(5.9)

add_text(s, Inches(0.7), col_y, col_w, Inches(0.35),
         "Why now", size=15, bold=True, color=GW_NAVY)
add_bullets(s, Inches(0.7), col_y + Inches(0.4), col_w, col_h - Inches(0.4), [
    "Q-Commerce rider workforce tripled (2022–2025), growing 34% YoY.",
    "IRDAI 2024 sandbox explicitly permits parametric micro-insurance pilots.",
    "Climate disruption in tier-1 cities up 2.4× vs. 10-yr baseline.",
    "UPI + Razorpay + Twilio Verify collapse the unit-economics floor.",
], size=11)

add_text(s, Inches(7.0), col_y, col_w, Inches(0.35),
         "Why incumbents can't", size=15, bold=True, color=GW_NAVY)
add_bullets(s, Inches(7.0), col_y + Inches(0.4), col_w, col_h - Inches(0.4), [
    "Legacy insurers need an adjuster — eats premium in one phone call.",
    "Risk is hyperlocal, hour-granular; actuarial tables don't bend that far.",
    "No distribution to gig workers without payslips or bank statements.",
    "Zero major insurers ship in 11 Indian languages.",
], size=11)

# =========================================================================
# SLIDE 4 — THE INSIGHT
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_slide_chrome(s, 4, "03 / Insight")
add_eyebrow(s, Inches(0.85), "The Reframe")

# Title with strikethrough
add_runs(s, Inches(0.7), Inches(1.15), Inches(12.3), Inches(1.8),
         [('This is not a ',     {'size': 32, 'bold': True, 'color': GW_NAVY}),
          ('claims',              {'size': 32, 'bold': True, 'color': RGBColor(0xB0, 0xB8, 0xC2), 'italic': True}),
          (' problem.\n',         {'size': 32, 'bold': True, 'color': GW_NAVY}),
          ('It is a ',            {'size': 32, 'bold': True, 'color': GW_NAVY}),
          ('parameters',          {'size': 32, 'bold': True, 'color': GW_ORANGE}),
          (' problem.',           {'size': 32, 'bold': True, 'color': GW_NAVY})],
         line_spacing=1.08)

# Two side-by-side cards
card_y = Inches(3.3); card_w = Inches(5.9); card_h = Inches(3.0)

# Left (legacy) — navy
add_rect(s, Inches(0.7), card_y, card_w, Inches(0.07), GW_ORANGE)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), card_y + Inches(0.07), card_w, card_h - Inches(0.07))
bg.fill.solid(); bg.fill.fore_color.rgb = GW_NAVY; bg.line.fill.background(); bg.shadow.inherit = False
add_pill(s, Inches(0.9), card_y + Inches(0.22), "LEGACY APPROACH",
         RGBColor(0x3A, 0x5E, 0x90), GW_WHITE, width=Inches(1.6))
add_text(s, Inches(0.9), card_y + Inches(0.65), card_w - Inches(0.4), Inches(0.4),
         "Indemnity insurance", size=18, bold=True, color=GW_WHITE)
add_bullets(s, Inches(0.9), card_y + Inches(1.1), card_w - Inches(0.4), card_h - Inches(1.2), [
    "Worker files claim with paperwork.",
    "Adjuster visits / verifies.",
    "42–75 day settlement.",
    "Unit economics impossible below ₹500/mo premium.",
], size=12, color=RGBColor(0xE4, 0xE8, 0xF0))

# Right (ours) — white with green top
add_rect(s, Inches(6.75), card_y, card_w, Inches(0.07), GW_GREEN)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.75), card_y + Inches(0.07), card_w, card_h - Inches(0.07))
bg.fill.solid(); bg.fill.fore_color.rgb = GW_WHITE; bg.line.color.rgb = RGBColor(0xD0, 0xD8, 0xE4); bg.line.width = Pt(0.5); bg.shadow.inherit = False
add_pill(s, Inches(6.95), card_y + Inches(0.22), "GIGSHIELD APPROACH",
         RGBColor(0xE6, 0xF7, 0xEC), GW_GREEN, width=Inches(1.8))
add_text(s, Inches(6.95), card_y + Inches(0.65), card_w - Inches(0.4), Inches(0.4),
         "Parametric insurance", size=18, bold=True, color=GW_NAVY)
add_bullets(s, Inches(6.95), card_y + Inches(1.1), card_w - Inches(0.4), card_h - Inches(1.2), [
    "Trigger fires on public data (weather/AQI/community).",
    "Claim auto-generated, not submitted.",
    "BAS fraud check + real Razorpay payout in < 10 min.",
    "Works at ₹15–₹199/week — human is out of the loop.",
], size=12)

# Math formula at bottom
add_rect(s, Inches(0.7), Inches(6.55), Inches(0.08), Inches(0.4), GW_ORANGE)
add_text(s, Inches(0.95), Inches(6.55), Inches(12), Inches(0.4),
         "Payout  =  min( cap_policy ,  hours_lost  ×  r_hourly  ×  σ_severity )    with σ ∈ [1.0, 1.8]",
         size=14, italic=True, bold=True, color=GW_NAVY, line_spacing=1.2)

# =========================================================================
# SLIDE 5 — THE SOLUTION (4-step flow)
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_slide_chrome(s, 5, "04 / Solution")
add_eyebrow(s, Inches(0.85), "How a disruption becomes a payout")
add_title(s, Inches(1.15), "Four steps. Zero paperwork. Under 10 minutes.",
          accent_word="Under 10 minutes.")

# 4 step boxes with arrows
flow_y = Inches(2.8)
step_w = Inches(2.7); step_h = Inches(1.9); gap = Inches(0.3)
steps = [
    ("1", "Parametric Trigger", "Open-Meteo, AQICN or ≥3 community reports cross threshold for a geo-fenced zone."),
    ("2", "Claim Auto-Created", "Claim Service finds every active policy in the zone, mints a claim per worker."),
    ("3", "11-signal BAS",      "Fraud Service scores every claim. < 0.30 auto-approves."),
    ("4", "Razorpay Payout",    "Real Orders API returns order_XXX. Worker sees confetti + UTR."),
]
x = Inches(0.7)
for i, (num, title, desc) in enumerate(steps):
    # Box
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, flow_y, step_w, step_h)
    box.fill.solid(); box.fill.fore_color.rgb = GW_WHITE
    box.line.color.rgb = GW_NAVY; box.line.width = Pt(2)
    box.shadow.inherit = False
    # Circle with number
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 x + step_w/2 - Inches(0.25), flow_y + Inches(0.18),
                                 Inches(0.5), Inches(0.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = GW_ORANGE; circle.line.fill.background()
    add_text(s, x + step_w/2 - Inches(0.25), flow_y + Inches(0.24),
             Inches(0.5), Inches(0.4),
             num, size=16, bold=True, color=GW_WHITE, align=PP_ALIGN.CENTER)
    # Title
    add_text(s, x + Inches(0.15), flow_y + Inches(0.8),
             step_w - Inches(0.3), Inches(0.35),
             title, size=12, bold=True, color=GW_NAVY, align=PP_ALIGN.CENTER)
    # Desc
    add_text(s, x + Inches(0.2), flow_y + Inches(1.15),
             step_w - Inches(0.4), step_h - Inches(1.15),
             desc, size=9, color=GW_GRAY1, align=PP_ALIGN.CENTER, line_spacing=1.3)
    # Arrow between
    if i < 3:
        arr_x = x + step_w + Inches(0.03)
        add_text(s, arr_x, flow_y + step_h/2 - Inches(0.2),
                 gap - Inches(0.05), Inches(0.4),
                 "›", size=28, bold=True, color=GW_ORANGE, align=PP_ALIGN.CENTER)
    x += step_w + gap

# Bottom 4 mini-stat cards
mini_y = Inches(5.1); mini_w = Inches(2.85); mini_h = Inches(1.6)
mini = [
    ("5 triggers",    "Heat >45°C, Rain >65 mm/h, AQI >400, Floods, Curfew."),
    ("11 signals",    "Behavioral anti-spoofing, ring-fraud aware."),
    ("3 gateways",    "Razorpay (real) / UPI / Stripe. VPA-routed."),
    ("11 languages",  "English + 10 Indian languages, auto-detect + persist."),
]
for i, (t, d) in enumerate(mini):
    left = Inches(0.7) + i * (mini_w + Inches(0.15))
    add_rect(s, left, mini_y, mini_w, Inches(0.07), GW_ORANGE)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, mini_y + Inches(0.07), mini_w, mini_h - Inches(0.07))
    bg.fill.solid(); bg.fill.fore_color.rgb = GW_WHITE
    bg.line.color.rgb = RGBColor(0xE0, 0xE6, 0xEE); bg.line.width = Pt(0.5)
    bg.shadow.inherit = False
    add_pill(s, left + Inches(0.15), mini_y + Inches(0.2), "LIVE",
             RGBColor(0xE6, 0xF7, 0xEC), GW_GREEN, width=Inches(0.7))
    add_text(s, left + Inches(0.15), mini_y + Inches(0.55), mini_w - Inches(0.3), Inches(0.35),
             t, size=13, bold=True, color=GW_NAVY)
    add_text(s, left + Inches(0.15), mini_y + Inches(0.95), mini_w - Inches(0.3), mini_h - Inches(1.0),
             d, size=9, color=GW_GRAY1, line_spacing=1.35)

# =========================================================================
# SLIDE 6 — LIVE PRODUCT
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_slide_chrome(s, 6, "05 / Demo")
add_eyebrow(s, Inches(0.85), "Try it right now")
add_title(s, Inches(1.15), "Deployed, monitored, and taking real payments.",
          accent_word="and taking real payments.")

# 3 portal cards
p_y = Inches(2.7); p_w = Inches(4.0); p_h = Inches(2.3)
portals = [
    ("gigshield.in",              "Main worker PWA. Real Twilio Verify OTP.",            "https://gigshield.in",                   GW_GREEN,  "ONLINE"),
    ("gigshield.in/admin/login",  "Live Ops SSE feed, fraud review, GPT-5.4 insights.",  "admin@gigshield.in · admin123",          GW_ORANGE, "ADMIN PANEL"),
    ("gigshield.in/demo",         "3 anti-spoofing scenarios: legit, spoof, ring.",      "No login required",                      GW_BLUE,   "SIMULATOR"),
]
for i, (title, desc, url, accent, pill) in enumerate(portals):
    left = Inches(0.7) + i * (p_w + Inches(0.2))
    add_rect(s, left, p_y, p_w, Inches(0.07), accent)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, p_y + Inches(0.07), p_w, p_h - Inches(0.07))
    bg.fill.solid(); bg.fill.fore_color.rgb = GW_WHITE
    bg.line.color.rgb = RGBColor(0xE0, 0xE6, 0xEE); bg.line.width = Pt(0.5)
    bg.shadow.inherit = False
    pill_bg = RGBColor(0xE6, 0xF7, 0xEC) if accent == GW_GREEN else \
              RGBColor(0xFE, 0xEE, 0xE0) if accent == GW_ORANGE else \
              RGBColor(0xE0, 0xF4, 0xFC)
    add_pill(s, left + Inches(0.2), p_y + Inches(0.25), pill, pill_bg, accent, width=Inches(1.3))
    add_text(s, left + Inches(0.2), p_y + Inches(0.62), p_w - Inches(0.4), Inches(0.4),
             title, size=15, bold=True, color=GW_NAVY)
    add_text(s, left + Inches(0.2), p_y + Inches(1.05), p_w - Inches(0.4), Inches(0.8),
             desc, size=11, color=GW_GRAY1, line_spacing=1.35)
    add_text(s, left + Inches(0.2), p_y + Inches(1.9), p_w - Inches(0.4), Inches(0.3),
             url, size=10, bold=True, color=GW_NAVY, font='Menlo')

# Razorpay proof callout
proof_y = Inches(5.3)
add_rect(s, Inches(0.7), proof_y, Inches(12), Inches(0.07), GW_GREEN)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), proof_y + Inches(0.07), Inches(12), Inches(1.6))
bg.fill.solid(); bg.fill.fore_color.rgb = GW_NAVY
bg.line.fill.background(); bg.shadow.inherit = False
add_text(s, Inches(0.95), proof_y + Inches(0.2), Inches(11), Inches(0.4),
         "Proof of real Razorpay integration",
         size=15, bold=True, color=GW_WHITE)
add_text(s, Inches(0.95), proof_y + Inches(0.65), Inches(11.5), Inches(1.1),
         'Production container just created a genuine order against Razorpay\'s live API:\n'
         '   {"id": "order_SeeezClhAFfttw", "entity": "order", "amount": 10000, "status": "created", "created_at": 1776449620}',
         size=11, color=RGBColor(0xD9, 0xE7, 0xFF), font='Menlo', line_spacing=1.4)

# =========================================================================
# SLIDE 7 — ARCHITECTURE
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_slide_chrome(s, 7, "06 / Architecture")
add_eyebrow(s, Inches(0.85), "End-to-end, as deployed")
add_title(s, Inches(1.15), "Production architecture. Single-region, horizontally scalable.",
          accent_word="Single-region, horizontally scalable.")

# ASCII arch diagram in a panel
arch = ("                              ┌────────────────────────────────┐\n"
        "                              │    Cloudflare  (WAF · SSL)      │\n"
        "                              └────────────────┬───────────────┘\n"
        "                                               │\n"
        "         ┌─────────────────────────────────────┴────────────────────────────┐\n"
        "         │           Hostinger VPS  (Ubuntu 22.04 · Docker)                 │\n"
        "         │                                                                   │\n"
        "         │   Nginx :443 ──────┐                                               │\n"
        "         │       │            │                                               │\n"
        "         │       ▼            ▼                                               │\n"
        "         │   ┌────────┐   ┌────────────────────────┐                          │\n"
        "         │   │React 19│   │ Node 20 · Express 4    │ ◀── api.razorpay.com     │\n"
        "         │   │  PWA   │   │ 15 Svcs · 9 Migrations │ ◀── Twilio Verify        │\n"
        "         │   │ 11 lang│   │ SSE /events/stream     │ ◀── OpenAI GPT-5.4       │\n"
        "         │   └────────┘   └──────┬─────────────────┘ ◀── AQICN / Open-Meteo   │\n"
        "         │                       │                                             │\n"
        "         │                       ▼                                             │\n"
        "         │           ┌──────────────┐   ┌──────────────┐                       │\n"
        "         │           │ PostgreSQL 16│   │  Redis 7      │                       │\n"
        "         │           │ 9 migrations │   │ AI+SSE cache  │                       │\n"
        "         │           │ PostGIS      │   │ weather 30m   │                       │\n"
        "         │           └──────────────┘   └──────────────┘                       │\n"
        "         └───────────────────────────────────────────────────────────────────┘")
arch_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.1),
                               Inches(12), Inches(3.2))
arch_box.fill.solid(); arch_box.fill.fore_color.rgb = GW_BG
arch_box.line.color.rgb = RGBColor(0xD0, 0xD8, 0xE4); arch_box.line.width = Pt(0.5)
arch_box.shadow.inherit = False
add_text(s, Inches(0.9), Inches(2.25), Inches(11.7), Inches(3.0),
         arch, size=9, color=GW_INK, font='Menlo', line_spacing=1.15)

# 4 KPI cards at bottom
b_y = Inches(5.5); b_w = Inches(2.85); b_h = Inches(1.5)
bottom = [
    ("15",    "",  "Backend services"),
    ("9",     "",  "Database migrations"),
    ("22",    "",  "Frontend pages"),
    ("5,580", "",  "Real dark stores"),
]
for i, (kpi, unit, label) in enumerate(bottom):
    add_kpi_card(s, Inches(0.7) + i * (b_w + Inches(0.15)), b_y, b_w, b_h,
                 kpi=kpi, unit=unit, label=label)

# =========================================================================
# SLIDE 8 — BAS ENGINE
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_slide_chrome(s, 8, "07 / Anti-Spoofing")
add_eyebrow(s, Inches(0.85), "Behavioral Anti-Spoofing · 11 Signals")
add_title(s, Inches(1.15), "Fraud at scale is boring, not theatrical.",
          accent_word="boring,")
add_subtitle(s, Inches(2.2),
    "The real attack isn't a GPS spoof — it's 50 riders sharing one device fingerprint "
    "from one hostel. BAS catches clusters, not just individuals.")

# Formula
add_rect(s, Inches(0.7), Inches(2.95), Inches(0.08), Inches(0.5), GW_ORANGE)
add_text(s, Inches(1.0), Inches(3.0), Inches(12), Inches(0.5),
         "BAS(c) = Σⱼ₌₁..₁₁  wⱼ · fⱼ(c),    Σ wⱼ = 1,    fⱼ(c) ∈ [0, 1]",
         size=14, italic=True, bold=True, color=GW_NAVY)

# Table header
t_y = Inches(3.65)
hdr = add_rect(s, Inches(0.7), t_y, Inches(12), Inches(0.35), GW_NAVY)
add_text(s, Inches(0.85), t_y + Inches(0.05), Inches(0.4), Inches(0.3),
         "#", size=10, bold=True, color=GW_WHITE)
add_text(s, Inches(1.3), t_y + Inches(0.05), Inches(3.2), Inches(0.3),
         "SIGNAL", size=10, bold=True, color=GW_WHITE)
add_text(s, Inches(4.6), t_y + Inches(0.05), Inches(0.8), Inches(0.3),
         "WEIGHT", size=10, bold=True, color=GW_WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(5.5), t_y + Inches(0.05), Inches(7.0), Inches(0.3),
         "WHAT IT CATCHES", size=10, bold=True, color=GW_WHITE)

rows = [
    ("1",  "IMPOSSIBLE_TRAVEL",        "0.18", "Haversine ÷ Δt > 120 km/h"),
    ("2",  "DEVICE_FINGERPRINT_REUSE", "0.15", "Same FP on ≥3 worker IDs"),
    ("3",  "VPA_COLLISION",            "0.12", "Same UPI VPA on ≥2 profiles"),
    ("4",  "CLAIM_BURST_RATE",         "0.10", "> 3 claims in rolling 24h"),
    ("5",  "GPS_JITTER_OUTLIER",       "0.10", "Lat/lng std-dev anomalous"),
    ("6",  "ZONE_MISMATCH",            "0.09", "Claim outside insured zones"),
    ("7",  "TIMESTAMP_ANOMALY",        "0.08", "Client clock skew > 5 min"),
    ("8",  "IP_ASN_SHARING",           "0.07", "≥5 workers on one ASN in 1h"),
    ("9",  "NEW_ACCOUNT_VELOCITY",     "0.05", "Account <72h old, ≥2 claims"),
    ("10", "PAYOUT_ACCOUNT_CHURN",     "0.04", "Beneficiary VPA changed in 7d"),
    ("11", "PATTERN_SIMILARITY",       "0.02", "Cosine sim to known ring"),
]
row_h = Inches(0.25)
for i, (n, sig, w, what) in enumerate(rows):
    ry = t_y + Inches(0.35) + i * row_h
    # Alternating background
    if i % 2 == 0:
        add_rect(s, Inches(0.7), ry, Inches(12), row_h, RGBColor(0xF5, 0xF7, 0xFA))
    add_text(s, Inches(0.85), ry + Inches(0.04), Inches(0.4), row_h,
             n, size=10, color=GW_INK)
    add_text(s, Inches(1.3), ry + Inches(0.04), Inches(3.2), row_h,
             sig, size=10, bold=True, color=GW_NAVY, font='Menlo')
    add_text(s, Inches(4.6), ry + Inches(0.04), Inches(0.8), row_h,
             w, size=10, bold=True, color=GW_INK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(5.5), ry + Inches(0.04), Inches(7.0), row_h,
             what, size=10, color=GW_INK)

# =========================================================================
# SLIDE 9 — PAYMENT ROUTING
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_slide_chrome(s, 9, "08 / Payments")
add_eyebrow(s, Inches(0.85), "Real Razorpay + VPA-suffix routing")
add_title(s, Inches(1.15), "Real API calls. Real order_XXX IDs. Graceful fallback.",
          accent_word="Graceful fallback.")

# Left: routing code
left_x = Inches(0.7); col_w = Inches(5.95)
add_text(s, left_x, Inches(2.3), col_w, Inches(0.4),
         "Gateway routing", size=15, bold=True, color=GW_NAVY)
# Code block bg
add_rect(s, left_x, Inches(2.75), col_w, Inches(1.6), GW_NAVY)
add_rect(s, left_x, Inches(2.75), Inches(0.05), Inches(1.6), GW_ORANGE)
routing_code = ("// VPA suffix → gateway\n"
                "*@paytm                      → Razorpay API  (REAL)\n"
                "*@upi, *@ybl, *@okhdfcbank   → UPI gateway   (mock)\n"
                "card, *@stripe               → Stripe         (mock)\n"
                "*@oksbi, *@ibl               → UPI gateway   (mock)\n"
                "fallback                     → Razorpay API  (REAL)")
add_text(s, left_x + Inches(0.2), Inches(2.85), col_w - Inches(0.3), Inches(1.5),
         routing_code, size=10, color=RGBColor(0xD9, 0xE7, 0xFF), font='Menlo', line_spacing=1.4)

add_text(s, left_x, Inches(4.55), col_w, Inches(0.4),
         "Why Orders API, not Payouts?", size=14, bold=True, color=GW_NAVY)
add_bullets(s, left_x, Inches(4.95), col_w, Inches(2.0), [
    "RazorpayX Payouts needs account activation — test accounts don't get it.",
    "Orders API is universal; returns real order_XXX any judge can look up.",
    "Settlement modeled out-of-band with realistic 1.8–3.2s delay.",
], size=11)

# Right: live API call code
right_x = Inches(6.95)
add_text(s, right_x, Inches(2.3), col_w, Inches(0.4),
         "Live API call flow", size=15, bold=True, color=GW_NAVY)
add_rect(s, right_x, Inches(2.75), col_w, Inches(3.0), GW_NAVY)
add_rect(s, right_x, Inches(2.75), Inches(0.05), Inches(3.0), GW_ORANGE)
api_code = ('POST https://api.razorpay.com/v1/orders\n'
            'Authorization: Basic base64(KEY_ID:SECRET)\n'
            'Content-Type: application/json\n'
            '\n'
            '{\n'
            '  "amount": 32000,     // paise\n'
            '  "currency": "INR",\n'
            '  "receipt": claimId,\n'
            '  "notes": {\n'
            '    "beneficiary_upi": "rider@paytm",\n'
            '    "claim_id": "CLM-ABC12"\n'
            '  }\n'
            '}\n'
            '// → order_SeeezClhAFfttw\n'
            '// → stored as transaction_ref')
add_text(s, right_x + Inches(0.2), Inches(2.85), col_w - Inches(0.3), Inches(2.9),
         api_code, size=9, color=RGBColor(0xD9, 0xE7, 0xFF), font='Menlo', line_spacing=1.35)

# Bottom callout (full width)
cb_y = Inches(6.15)
add_rect(s, Inches(0.7), cb_y, Inches(12), Inches(0.07), GW_GREEN)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), cb_y + Inches(0.07), Inches(12), Inches(0.85))
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xE6, 0xF7, 0xEC)
bg.line.fill.background(); bg.shadow.inherit = False
add_pill(s, Inches(0.9), cb_y + Inches(0.2), "VERIFIED IN PRODUCTION",
         GW_GREEN, GW_WHITE, width=Inches(2.0))
add_text(s, Inches(3.0), cb_y + Inches(0.2), Inches(9.5), Inches(0.6),
         '18:13 UTC, April 17: prod container created order_SeeezClhAFfttw against Razorpay\'s live API — '
         'logged, audited, grep-able in the dashboard.',
         size=11, bold=True, color=GW_NAVY, line_spacing=1.3)

# =========================================================================
# SLIDE 10 — AI LAYER
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_slide_chrome(s, 10, "09 / AI")
add_eyebrow(s, Inches(0.85), "OpenAI GPT-5.4 · Three use cases")
add_title(s, Inches(1.15), "Natural-language reasoning where humans actually read.",
          accent_word="actually read.")
add_subtitle(s, Inches(2.2),
    "No Python microservice. Pure TypeScript via the official OpenAI SDK. Redis-cached.")

# 3 cards
c_y = Inches(3.1); c_w = Inches(4.0); c_h = Inches(3.0)
cards = [
    ("WORKER-FACING", GW_ORANGE,
     "Risk Narrative",
     '"Your zone (Powai) has HIGH risk because it floods during monsoon peaks and your '
     'delivery hours overlap with peak heat..."',
     "generateRiskNarrative(worker, score)\ncache: Redis 1h TTL"),
    ("ADMIN-FACING",  GW_NAVY,
     "Claim Assessment",
     '"Claim shows IMPOSSIBLE_TRAVEL of 450 km/h between pings. Device FP reused from flagged '
     'account wr_abc123. Recommend REJECT."',
     "generateClaimAssessment(claim, bas)\ncontext: full BAS breakdown"),
    ("DASHBOARD",     GW_BLUE,
     "Fraud Pattern Summary",
     '"3 claims flagged this week. Pattern: shared device fingerprints from Koramangala suggest '
     'coordinated ring fraud."',
     "generateFraudSummary(recent)\ncache: Redis 10 min TTL"),
]
for i, (pill_txt, accent, title, body, tech) in enumerate(cards):
    left = Inches(0.7) + i * (c_w + Inches(0.2))
    add_rect(s, left, c_y, c_w, Inches(0.07), accent)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, c_y + Inches(0.07), c_w, c_h - Inches(0.07))
    bg.fill.solid(); bg.fill.fore_color.rgb = GW_WHITE
    bg.line.color.rgb = RGBColor(0xE0, 0xE6, 0xEE); bg.line.width = Pt(0.5)
    bg.shadow.inherit = False
    pill_bg = RGBColor(0xFE, 0xEE, 0xE0) if accent == GW_ORANGE else \
              RGBColor(0xE0, 0xE8, 0xF2) if accent == GW_NAVY else \
              RGBColor(0xE0, 0xF4, 0xFC)
    add_pill(s, left + Inches(0.2), c_y + Inches(0.22), pill_txt, pill_bg, accent, width=Inches(1.5))
    add_text(s, left + Inches(0.2), c_y + Inches(0.6), c_w - Inches(0.4), Inches(0.4),
             title, size=15, bold=True, color=GW_NAVY)
    add_text(s, left + Inches(0.2), c_y + Inches(1.05), c_w - Inches(0.4), Inches(1.5),
             body, size=10, color=GW_INK, italic=True, line_spacing=1.4)
    add_text(s, left + Inches(0.2), c_y + Inches(2.35), c_w - Inches(0.4), Inches(0.6),
             tech, size=9, color=GW_GRAY1, font='Menlo', line_spacing=1.4)

# Bottom full-width navy callout
rag_y = Inches(6.25)
add_rect(s, Inches(0.7), rag_y, Inches(12), Inches(0.07), GW_ORANGE)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), rag_y + Inches(0.07), Inches(12), Inches(0.78))
bg.fill.solid(); bg.fill.fore_color.rgb = GW_NAVY
bg.line.fill.background(); bg.shadow.inherit = False
add_text(s, Inches(0.9), rag_y + Inches(0.15), Inches(11), Inches(0.35),
         "RAG context we inject", size=12, bold=True, color=GW_WHITE)
add_text(s, Inches(0.9), rag_y + Inches(0.47), Inches(11.5), Inches(0.4),
         "Worker BAS breakdown · claim history · zone weather history · device fingerprint · "
         "ring-fraud signal output · policy metadata — structured as typed JSON so GPT-5.4 "
         "returns bounded, citable reasoning.",
         size=10, color=RGBColor(0xD9, 0xE7, 0xFF), line_spacing=1.3)

# =========================================================================
# SLIDE 11 — 11 LANGUAGES
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_slide_chrome(s, 11, "10 / i18n")
add_eyebrow(s, Inches(0.85), "Vernacular is not a checkbox")
add_title(s, Inches(1.15), "11 Indian languages. Native scripts. Auto-detected.",
          accent_word="Native scripts.")
add_subtitle(s, Inches(2.2),
    'First three test riders asked "Tamil me hai kya?" within 2 minutes. We shipped i18n that evening.')

# 12 language cards (11 + stack)
langs = [
    ("English",  "en · Latin"),
    ("हिन्दी",      "hi · Devanagari"),
    ("தமிழ்",    "ta · Tamil"),
    ("తెలుగు",   "te · Telugu"),
    ("বাংলা",      "bn · Bengali"),
    ("मराठी",     "mr · Devanagari"),
    ("ಕನ್ನಡ",    "kn · Kannada"),
    ("ગુજરાતી",   "gu · Gujarati"),
    ("ਪੰਜਾਬੀ",   "pa · Gurmukhi"),
    ("മലയാളം",  "ml · Malayalam"),
    ("ଓଡ଼ିଆ",     "or · Odia"),
]
g_y = Inches(3.0); g_w = Inches(2.0); g_h = Inches(1.05)
cols = 6
for i, (lang, meta) in enumerate(langs):
    col = i % cols; row = i // cols
    left = Inches(0.7) + col * (g_w + Inches(0.05))
    top  = g_y + row * (g_h + Inches(0.15))
    accent = GW_ORANGE if lang == "English" else RGBColor(0xE0, 0xE6, 0xEE)
    add_rect(s, left, top, g_w, Inches(0.05), accent)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.05), g_w, g_h - Inches(0.05))
    bg.fill.solid(); bg.fill.fore_color.rgb = GW_WHITE
    bg.line.color.rgb = RGBColor(0xE0, 0xE6, 0xEE); bg.line.width = Pt(0.5)
    bg.shadow.inherit = False
    add_text(s, left, top + Inches(0.2), g_w, Inches(0.5),
             lang, size=20, bold=True, color=GW_NAVY, align=PP_ALIGN.CENTER)
    add_text(s, left, top + Inches(0.75), g_w, Inches(0.25),
             meta, size=8, color=GW_GRAY1, align=PP_ALIGN.CENTER)

# Stack card (12th slot)
col = 11 % cols; row = 11 // cols
left = Inches(0.7) + col * (g_w + Inches(0.05))
top  = g_y + row * (g_h + Inches(0.15))
add_rect(s, left, top, g_w, Inches(0.05), GW_ORANGE)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.05), g_w, g_h - Inches(0.05))
bg.fill.solid(); bg.fill.fore_color.rgb = GW_NAVY
bg.line.fill.background(); bg.shadow.inherit = False
add_text(s, left, top + Inches(0.14), g_w, Inches(0.3),
         "STACK", size=9, bold=True, color=RGBColor(0xB0, 0xC0, 0xD8), align=PP_ALIGN.CENTER)
add_text(s, left, top + Inches(0.38), g_w, Inches(0.7),
         "i18next\n+ react-i18next\n+ browser-detector", size=10, bold=True,
         color=GW_WHITE, align=PP_ALIGN.CENTER, line_spacing=1.25)

# Footer: implementation
imp_y = Inches(5.9)
add_text(s, Inches(0.7), imp_y, Inches(12), Inches(0.8),
         "Implementation: auto-detected from device locale on first load, persisted to localStorage, "
         "switchable via globe icon in every page header. 40+ keys per locale covering Welcome, Login, "
         "Dashboard, Profile and all CTAs. 440+ translated strings total.",
         size=11, color=GW_GRAY1, line_spacing=1.4)

# =========================================================================
# SLIDE 12 — LIVE OPS WAR ROOM
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_slide_chrome(s, 12, "11 / Real-Time")
add_eyebrow(s, Inches(0.85), "Server-Sent Events · < 800 ms latency")
add_title(s, Inches(1.15), "Static dashboards feel like schoolwork. Ours breathes.",
          accent_word="Ours breathes.")

# Left — 6 event chips
col_x = Inches(0.7); col_w = Inches(5.9)
add_text(s, col_x, Inches(2.65), col_w, Inches(0.4),
         "6 typed event streams", size=15, bold=True, color=GW_NAVY)
events = [
    ("RED",    GW_RED,    "TRIGGER_FIRED",     "new disruption detected"),
    ("AMBER",  GW_AMBER,  "CLAIM_CREATED",     "auto-claim minted"),
    ("ORANGE", GW_ORANGE, "FRAUD_CHECKED",     "BAS score computed"),
    ("GREEN",  GW_GREEN,  "PAYOUT_SENT",       "Razorpay accepted"),
    ("BLUE",   GW_BLUE,   "WORKER_REGISTERED", "rider onboarded"),
    ("PURPLE", RGBColor(0x93, 0x33, 0xEA), "POLICY_CREATED", "coverage purchased"),
]
ey = Inches(3.15)
row_h = Inches(0.45)
for i, (lbl, accent, name, desc) in enumerate(events):
    row_y = ey + i * (row_h + Inches(0.08))
    # Tint bg
    tint_r, tint_g, tint_b = accent[0], accent[1], accent[2]
    # use a pre-mix for simplicity — 0.08 alpha-ish
    tint = RGBColor(min(255, tint_r + (255 - tint_r) * 92 // 100),
                    min(255, tint_g + (255 - tint_g) * 92 // 100),
                    min(255, tint_b + (255 - tint_b) * 92 // 100))
    add_rect(s, col_x, row_y, col_w, row_h, tint)
    add_pill(s, col_x + Inches(0.15), row_y + Inches(0.09), lbl, accent, GW_WHITE, width=Inches(0.9))
    add_runs(s, col_x + Inches(1.15), row_y + Inches(0.12), col_w - Inches(1.3), row_h,
             [(name + '  ', {'size': 11, 'bold': True, 'color': GW_NAVY, 'font': 'Menlo'}),
              (' — ' + desc, {'size': 11, 'color': GW_INK})])

# Right — implementation code
right_x = Inches(6.95); right_w = Inches(5.9)
add_text(s, right_x, Inches(2.65), right_w, Inches(0.4),
         "Implementation", size=15, bold=True, color=GW_NAVY)
add_rect(s, right_x, Inches(3.1), right_w, Inches(3.1), GW_NAVY)
add_rect(s, right_x, Inches(3.1), Inches(0.05), Inches(3.1), GW_ORANGE)
code = ("// services/events.service.ts\n"
        "export class EventBus extends EventEmitter {\n"
        "  emit<T>(type: EventType, payload: T) {\n"
        "    super.emit('event', { type, payload, ts: Date.now() });\n"
        "  }\n"
        "}\n"
        "\n"
        "// routes/events.routes.ts\n"
        "router.get('/admin/events/stream', requireAdmin, (req, res) => {\n"
        "  res.writeHead(200, {\n"
        "    'Content-Type': 'text/event-stream',\n"
        "    'Cache-Control': 'no-cache',\n"
        "  });\n"
        "  bus.on('event', e => res.write(`data: ${JSON.stringify(e)}\\n\\n`));\n"
        "  // 30s heartbeat survives Cloudflare idle timeout\n"
        "  setInterval(() => res.write(': hb\\n\\n'), 30000);\n"
        "});")
add_text(s, right_x + Inches(0.2), Inches(3.2), right_w - Inches(0.3), Inches(3.0),
         code, size=9, color=RGBColor(0xD9, 0xE7, 0xFF), font='Menlo', line_spacing=1.35)

# Callout
cout_y = Inches(6.35)
add_rect(s, right_x, cout_y, right_w, Inches(0.07), GW_GREEN)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, cout_y + Inches(0.07), right_w, Inches(0.6))
bg.fill.solid(); bg.fill.fore_color.rgb = GW_WHITE
bg.line.color.rgb = RGBColor(0xE0, 0xE6, 0xEE); bg.line.width = Pt(0.5)
bg.shadow.inherit = False
add_text(s, right_x + Inches(0.2), cout_y + Inches(0.17), right_w - Inches(0.3), Inches(0.5),
         'Each event animates into the feed. Live payout counter ticks upward on every '
         'PAYOUT_SENT — our "stop scrolling" demo moment.',
         size=10, color=GW_INK, line_spacing=1.3)

# =========================================================================
# SLIDE 13 — COMMUNITY REPORTS (DIFFERENTIATOR)
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_slide_chrome(s, 13, "12 / Community")
add_eyebrow(s, Inches(0.85), "Workers as hyperlocal sensors")
add_title(s, Inches(1.15), "Weather APIs lag 15–30 minutes.\nA rider in water already knows.",
          accent_word="already knows.")

# Left — feedback loop
left_x = Inches(0.7); col_w = Inches(6.0)
add_text(s, left_x, Inches(3.0), col_w, Inches(0.4),
         "The feedback loop", size=15, bold=True, color=GW_NAVY)

# Orange vertical line
add_rect(s, left_x + Inches(0.22), Inches(3.55), Inches(0.02), Inches(3.3), GW_ORANGE)
loops = [
    ("1", GW_ORANGE, 'Worker taps "Report Condition"',
     "Icons: Heavy Rain / Flood / Heat / Pollution / Strike. GPS auto-captured, severity slider."),
    ("2", GW_ORANGE, "Cluster detection",
     "≥3 reports within 2 km + 15 min → provisional DisruptionEvent → claims begin."),
    ("3", GW_ORANGE, "API confirmation",
     "Later, when Open-Meteo/AQICN confirms the same event: boostTrustOnApiConfirm()."),
    ("4", GW_GREEN,  "Trust Score +50",
     "Every original reporter earns +50 Trust Score. Honest reporting gets rewarded."),
]
step_y = Inches(3.55)
for i, (num, color, title, desc) in enumerate(loops):
    y = step_y + i * Inches(0.85)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, left_x, y, Inches(0.5), Inches(0.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = color; circle.line.fill.background()
    add_text(s, left_x, y + Inches(0.05), Inches(0.5), Inches(0.4),
             num, size=13, bold=True, color=GW_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, left_x + Inches(0.7), y, col_w - Inches(0.8), Inches(0.3),
             title, size=12, bold=True, color=GW_NAVY)
    add_text(s, left_x + Inches(0.7), y + Inches(0.3), col_w - Inches(0.8), Inches(0.6),
             desc, size=10, color=GW_GRAY1, line_spacing=1.3)

# Right — navy "why this wins" card + on-disk proof
right_x = Inches(7.0); r_w = Inches(5.85)

# Navy card
card_y = Inches(3.0)
add_rect(s, right_x, card_y, r_w, Inches(0.07), GW_ORANGE)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, card_y + Inches(0.07), r_w, Inches(2.0))
bg.fill.solid(); bg.fill.fore_color.rgb = GW_NAVY
bg.line.fill.background(); bg.shadow.inherit = False
add_pill(s, right_x + Inches(0.2), card_y + Inches(0.25), "WHY THIS WINS",
         RGBColor(0x5E, 0x2E, 0x10), GW_ORANGE2, width=Inches(1.6))
add_text(s, right_x + Inches(0.2), card_y + Inches(0.65), r_w - Inches(0.4), Inches(0.45),
         "No other team ships this.", size=18, bold=True, color=GW_WHITE)
add_text(s, right_x + Inches(0.2), card_y + Inches(1.15), r_w - Inches(0.4), Inches(0.9),
         "We turn 1.2M riders into a distributed sensor network that outruns commercial weather "
         "APIs by minutes. Over time the system becomes more accurate than its upstream data.",
         size=11, color=RGBColor(0xE0, 0xE8, 0xF2), line_spacing=1.4)

# On-disk proof table
proof_y = Inches(5.2)
add_text(s, right_x, proof_y, r_w, Inches(0.3),
         "On-disk proof", size=13, bold=True, color=GW_NAVY)
proofs = [
    ("migrations/005_community_reports.ts", "table + trust_score column"),
    ("services/community.service.ts",       "checkCluster · boost"),
    ("routes/community.routes.ts",          "POST /report · GET /pending"),
    ("components/ReportConditionModal.tsx", "5-icon worker UI"),
]
rh = Inches(0.32)
for i, (path, desc) in enumerate(proofs):
    rowy = proof_y + Inches(0.35) + i * rh
    if i % 2 == 0:
        add_rect(s, right_x, rowy, r_w, rh, RGBColor(0xF5, 0xF7, 0xFA))
    add_text(s, right_x + Inches(0.1), rowy + Inches(0.05), Inches(3.2), rh,
             path, size=9, bold=True, color=GW_NAVY, font='Menlo')
    add_text(s, right_x + Inches(3.35), rowy + Inches(0.05), r_w - Inches(3.45), rh,
             desc, size=10, color=GW_INK)

# =========================================================================
# SLIDE 14 — UNIT ECONOMICS
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_slide_chrome(s, 14, "13 / Economics")
add_eyebrow(s, Inches(0.85), "Unit economics · Weekly premium")
add_title(s, Inches(1.15), "₹29–₹199/week. Loss ratio 62% projected.",
          accent_word="Loss ratio 62% projected.")

# Pricing table
t_y = Inches(2.6)
hdr = add_rect(s, Inches(0.7), t_y, Inches(12), Inches(0.35), GW_NAVY)
headers = [("TIER", 0.9), ("COVERAGE", 2.5), ("PREMIUM /WK", 1.7), ("MAX PAYOUT /EVENT", 2.0),
           ("EXP. CLAIMS /YR", 1.7), ("BURN", 0.9)]
x = Inches(0.85)
for h, w in headers:
    add_text(s, x, t_y + Inches(0.05), Inches(w), Inches(0.3),
             h, size=10, bold=True, color=GW_WHITE)
    x += Inches(w + 0.1)

tiers = [
    ("Basic",    GW_NAVY,   "40% of lost income", "₹29 – ₹70",   "₹400",   "1.8", "58%"),
    ("Standard", GW_BLUE,   "60% of lost income", "₹71 – ₹140",  "₹750",   "2.4", "62%"),
    ("Premium",  GW_ORANGE, "80% of lost income", "₹141 – ₹199", "₹1,200", "3.1", "64%"),
]
rh = Inches(0.55)
for i, (tier, tier_color, cov, prem, payout, claims, burn) in enumerate(tiers):
    rowy = t_y + Inches(0.35) + i * rh
    if i % 2 == 0:
        add_rect(s, Inches(0.7), rowy, Inches(12), rh, RGBColor(0xF5, 0xF7, 0xFA))
    add_pill(s, Inches(0.85), rowy + Inches(0.15), tier,
             RGBColor(tier_color[0], tier_color[1], tier_color[2]) if False else
             (RGBColor(0xE0, 0xE8, 0xF2) if tier_color == GW_NAVY else
              RGBColor(0xE0, 0xF4, 0xFC) if tier_color == GW_BLUE else
              RGBColor(0xFE, 0xEE, 0xE0)),
             tier_color, width=Inches(0.9))
    x = Inches(1.85)
    for val, w in [(cov, 2.5), (prem, 1.7), (payout, 2.0), (claims, 1.7), (burn, 0.9)]:
        add_text(s, x, rowy + Inches(0.18), Inches(w), Inches(0.3),
                 val, size=11, color=GW_INK)
        x += Inches(w + 0.1)

# 4 KPIs
k_y = Inches(5.2); k_w = Inches(2.85); k_h = Inches(1.9)
kpis = [
    ("62", "%",  "Target loss ratio",  "Industry bench: 58–72% for parametric micro-insurance.", False),
    ("~0", "%",  "Adjuster cost",      "Fully automated. Human only on UNDER_REVIEW.",           False),
    ("18", "%",  "Gross margin",       "After fees, GST, infra and ~8% BAS-flagged leak.",        False),
    ("208", "Cr","Year-3 GWP target",  "At 15% SOM of Q-Comm riders (~180K policies).",           True),
]
for i, (v, u, lbl, desc, navy) in enumerate(kpis):
    add_kpi_card(s, Inches(0.7) + i * (k_w + Inches(0.15)), k_y, k_w, k_h,
                 kpi=v, unit=u, label=lbl, desc=desc, navy=navy,
                 accent=GW_ORANGE if not navy else GW_ORANGE)

# =========================================================================
# SLIDE 15 — TECH STACK
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_slide_chrome(s, 15, "14 / Tech Stack")
add_eyebrow(s, Inches(0.85), "What's actually in the repo")
add_title(s, Inches(1.15), "No vapourware. Every claim grep-able.",
          accent_word="Every claim grep-able.")

cats = [
    ("Frontend", [
        "React 19 + TypeScript",
        "Vite, Tailwind CSS v4",
        "Zustand for state",
        "i18next (11 locales)",
        "sonner + canvas-confetti",
        "react-countup, jsPDF",
        "@react-google-maps/api",
    ]),
    ("Backend", [
        "Node 20 + Express 4",
        "Knex (9 migrations)",
        "PostgreSQL 16 + PostGIS",
        "Redis 7 (cache + SSE)",
        "Zod, Pino, node-cron",
        "JWT + refresh tokens",
        "15 service modules",
    ]),
    ("Integrations", [
        "Razorpay Orders API (real)",
        "Twilio Verify (real OTP)",
        "OpenAI GPT-5.4",
        "Open-Meteo (free)",
        "AQICN (air quality)",
        "Google Maps JS API",
    ]),
    ("Infra", [
        "Hostinger VPS · Ubuntu 22",
        "Docker + Compose (prod)",
        "Nginx reverse proxy",
        "Cloudflare (SSL, WAF)",
        "Multi-stage Dockerfile",
        "Image < 300 MB",
    ]),
    ("Quality", [
        "TypeScript strict",
        "Zod env validation",
        "Hash-chained audit (mig 006)",
        "Error boundaries (React)",
        "Skeleton loaders",
        "Rate-limit + Helmet",
    ]),
    ("Data", [
        "~5,580 real dark stores",
        "11 i18n locale JSONs",
        "9 migrations + 4 seeds",
        "22 frontend pages",
        "11-signal BAS engine",
        "3 payment gateways",
    ]),
]
cw = Inches(4.0); ch = Inches(2.05); gap = Inches(0.1)
start_y = Inches(2.7)
for idx, (cat, items) in enumerate(cats):
    col = idx % 3; row = idx // 3
    left = Inches(0.7) + col * (cw + gap)
    top  = start_y + row * (ch + Inches(0.2))
    add_rect(s, left, top, cw, Inches(0.07), GW_ORANGE)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.07), cw, ch - Inches(0.07))
    bg.fill.solid(); bg.fill.fore_color.rgb = GW_WHITE
    bg.line.color.rgb = RGBColor(0xE0, 0xE6, 0xEE); bg.line.width = Pt(0.5)
    bg.shadow.inherit = False
    add_text(s, left + Inches(0.2), top + Inches(0.2), cw - Inches(0.4), Inches(0.3),
             cat, size=14, bold=True, color=GW_NAVY)
    add_bullets(s, left + Inches(0.2), top + Inches(0.55), cw - Inches(0.4), ch - Inches(0.6),
                items, size=10)

# =========================================================================
# SLIDE 16 — TRACTION / SCORECARD
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_slide_chrome(s, 16, "15 / Traction")
add_eyebrow(s, Inches(0.85), "What shipped in 13 days")
add_title(s, Inches(1.15), "From Phase-2 prototype to real Razorpay payouts.",
          accent_word="real Razorpay payouts.")

# 4 stats
sy = Inches(2.55); sw = Inches(2.95); sh = Inches(1.3)
stats = [
    ("100", "%",   "Phase-3 plan shipped"),
    ("order_SeeezClhAFfttw", "", "Live Razorpay production order"),
    ("<1",  "s",   "SSE event latency"),
    ("11",  "lang", "Indian languages · native scripts"),
]
for i, (big, unit, lbl) in enumerate(stats):
    left = Inches(0.7) + i * (sw + Inches(0.1))
    add_rect(s, left, sy, sw, Inches(0.07), GW_ORANGE)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, sy + Inches(0.07), sw, sh - Inches(0.07))
    bg.fill.solid(); bg.fill.fore_color.rgb = GW_WHITE
    bg.line.color.rgb = RGBColor(0xE0, 0xE6, 0xEE); bg.line.width = Pt(0.5)
    bg.shadow.inherit = False
    if i == 1:
        # Order id — use Menlo + smaller
        add_text(s, left + Inches(0.15), sy + Inches(0.35), sw - Inches(0.3), Inches(0.4),
                 big, size=14, bold=True, color=GW_NAVY, font='Menlo')
    else:
        add_runs(s, left + Inches(0.15), sy + Inches(0.25), sw - Inches(0.3), Inches(0.5),
                 [(big, {'size': 32, 'bold': True, 'color': GW_NAVY}),
                  (' ' + unit, {'size': 14, 'bold': True, 'color': GW_ORANGE})])
    add_text(s, left + Inches(0.15), sy + Inches(0.85), sw - Inches(0.3), Inches(0.4),
             lbl, size=10, bold=True, color=GW_GRAY1, line_spacing=1.3)

# Changelog table
add_text(s, Inches(0.7), Inches(4.15), Inches(12), Inches(0.4),
         "Phase 3 changelog", size=14, bold=True, color=GW_NAVY)
ct_y = Inches(4.55)
add_rect(s, Inches(0.7), ct_y, Inches(12), Inches(0.3), GW_NAVY)
add_text(s, Inches(0.85),  ct_y + Inches(0.04), Inches(0.4), Inches(0.25),
         "#",       size=9, bold=True, color=GW_WHITE)
add_text(s, Inches(1.3),   ct_y + Inches(0.04), Inches(5.2), Inches(0.25),
         "FEATURE", size=9, bold=True, color=GW_WHITE)
add_text(s, Inches(6.65),  ct_y + Inches(0.04), Inches(5.0), Inches(0.25),
         "FILES",   size=9, bold=True, color=GW_WHITE)
add_text(s, Inches(11.75), ct_y + Inches(0.04), Inches(0.9), Inches(0.25),
         "STATUS",  size=9, bold=True, color=GW_WHITE)

changelog = [
    ("1",  "Real Razorpay Orders API (with mock fallback)", "external/payments/razorpay.mock.ts"),
    ("2",  "VPA-suffix payment routing",                    "external/payments/index.ts"),
    ("3",  "Payout lifecycle (UTR, fee, tax, response)",    "migrations/009_payout_lifecycle.ts"),
    ("4",  "SSE Live Ops feed with 6 typed events",         "routes/events.routes.ts"),
    ("5",  "Community hyperlocal reports + cluster",         "services/community.service.ts"),
    ("6",  "Trust Score system",                            "migrations/005_community_reports.ts"),
    ("7",  "11-language i18n",                              "frontend/src/i18n/locales/"),
    ("8",  "Worker celebration (confetti + UTR toast)",     "sonner + canvas-confetti"),
    ("9",  "Hero stats on Welcome + public stats",          "components/HeroStats.tsx"),
    ("10", "Ring-fraud hardening (ASN + FP indexing)",      "migrations/008_ring_fraud_hardening.ts"),
]
rh = Inches(0.22)
for i, (n, feat, files) in enumerate(changelog):
    rowy = ct_y + Inches(0.3) + i * rh
    if i % 2 == 0:
        add_rect(s, Inches(0.7), rowy, Inches(12), rh, RGBColor(0xF5, 0xF7, 0xFA))
    add_text(s, Inches(0.85), rowy + Inches(0.03), Inches(0.4), rh,
             n, size=9, color=GW_INK)
    add_text(s, Inches(1.3),  rowy + Inches(0.03), Inches(5.2), rh,
             feat, size=9, color=GW_INK)
    add_text(s, Inches(6.65), rowy + Inches(0.03), Inches(5.0), rh,
             files, size=8, color=GW_NAVY, font='Menlo', bold=True)
    add_pill(s, Inches(11.75), rowy + Inches(0.01), "LIVE",
             RGBColor(0xE6, 0xF7, 0xEC), GW_GREEN, width=Inches(0.6))

# =========================================================================
# SLIDE 17 — ROADMAP / ASK
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_slide_chrome(s, 17, "16 / Roadmap")
add_eyebrow(s, Inches(0.85), "Path to 150,000 policies")
add_title(s, Inches(1.15), "Next 12 months. Scale + regulation + distribution.",
          accent_word="Scale + regulation + distribution.")

# 3 cards
cy = Inches(2.65); cw = Inches(4.0); ch = Inches(3.4)
phases = [
    ("Q2 2026", GW_BLUE, "Pilot with one Q-Comm platform", [
        "Blinkit or Zepto pilot — 5,000 riders in Mumbai",
        "RazorpayX activation for real payouts",
        "IRDAI sandbox application",
        "Integration with platform payroll",
    ], False),
    ("Q3–Q4 2026", GW_ORANGE, "Regulatory + scale", [
        "Partner with a licensed general insurer",
        "ML-based ring detection (graph embeddings)",
        "WhatsApp-native onboarding channel",
        "Expand to 3 cities · 50,000 riders",
    ], False),
    ("THE ASK", GW_ORANGE, "What Guidewire can unlock", [
        "Introduction to licensed insurance partners",
        "Actuarial expertise review of our reserves",
        "Product mentorship for the pilot launch",
        "Angel investment consideration",
    ], True),
]
for i, (pill_txt, accent, title, items, navy) in enumerate(phases):
    left = Inches(0.7) + i * (cw + Inches(0.2))
    add_rect(s, left, cy, cw, Inches(0.07), accent)
    bg_color = GW_NAVY if navy else GW_WHITE
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, cy + Inches(0.07), cw, ch - Inches(0.07))
    bg.fill.solid(); bg.fill.fore_color.rgb = bg_color
    bg.line.color.rgb = RGBColor(0xE0, 0xE6, 0xEE) if not navy else GW_NAVY
    bg.line.width = Pt(0.5)
    bg.shadow.inherit = False
    pill_bg = RGBColor(0x5E, 0x2E, 0x10) if navy else \
              (RGBColor(0xE0, 0xF4, 0xFC) if accent == GW_BLUE else RGBColor(0xFE, 0xEE, 0xE0))
    pill_fg = GW_ORANGE2 if navy else accent
    add_pill(s, left + Inches(0.2), cy + Inches(0.22), pill_txt, pill_bg, pill_fg, width=Inches(1.6))
    title_color = GW_WHITE if navy else GW_NAVY
    add_text(s, left + Inches(0.2), cy + Inches(0.6), cw - Inches(0.4), Inches(0.6),
             title, size=15, bold=True, color=title_color, line_spacing=1.25)
    item_color = RGBColor(0xE4, 0xE8, 0xF0) if navy else GW_INK
    add_bullets(s, left + Inches(0.2), cy + Inches(1.4), cw - Inches(0.4), ch - Inches(1.5),
                items, size=11, color=item_color)

# Big orange callout
cal_y = Inches(6.25)
add_rect(s, Inches(0.7), cal_y, Inches(12), Inches(0.7), GW_ORANGE)
add_text(s, Inches(0.9), cal_y + Inches(0.2), Inches(11.5), Inches(0.4),
         "We built all of this in 45 days.    Imagine what we ship in 12 months.",
         size=17, bold=True, color=GW_WHITE, align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 18 — THANK YOU / OUTRO
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, GW_NAVY)

# Orange corner
add_rect(s, Inches(11.5), 0, Inches(1.8), Inches(0.4), GW_ORANGE)
# Phase pill top-right
add_pill(s, Inches(9.5), Inches(0.6), "GUIDEWIRE DEVTRAILS 2026 · PHASE 3 FINAL",
         RGBColor(0x5E, 0x2E, 0x10), GW_ORANGE2, width=Inches(3.5))

# Big centered GigShield logo above the sign-off
_thx_logo_h = Inches(1.35)
_thx_logo = add_logo(s, Inches(0), Inches(1.1), height=_thx_logo_h, on_dark=True)
# Re-centre horizontally now that python-pptx has computed the picture width
_thx_logo.left = int((SLIDE_W - _thx_logo.width) / 2)

# Huge "Thank you."
add_runs(s, Inches(1), Inches(2.7), Inches(11.3), Inches(1.5),
         [('Thank ', {'size': 64, 'bold': True, 'color': GW_WHITE}),
          ('you',    {'size': 64, 'bold': True, 'color': GW_ORANGE}),
          ('.',      {'size': 64, 'bold': True, 'color': GW_WHITE})],
         align=PP_ALIGN.CENTER)

# Tagline
add_text(s, Inches(1), Inches(4.0), Inches(11.3), Inches(0.5),
         "Protecting India's gig workers, one disruption at a time.",
         size=17, color=RGBColor(0xD0, 0xD8, 0xE8), align=PP_ALIGN.CENTER, italic=True)

# 4 link cards
lc_y = Inches(5.0); lc_w = Inches(2.9); lc_h = Inches(1.3)
lnks = [
    ("LIVE PRODUCT",  "gigshield.in"),
    ("DEMO VIDEO",    "youtu.be/bx8AVAU_amk"),
    ("SOURCE CODE",   "github.com/uroy80/\nTeamVoid-Guidewire-Devtrails"),
    ("TEAM",          "Team Void"),
]
total_w = lc_w * 4 + Inches(0.15) * 3
start_x = (SLIDE_W - total_w) / 2
for i, (k, v) in enumerate(lnks):
    left = start_x + i * (lc_w + Inches(0.15))
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, lc_y, lc_w, lc_h)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1A, 0x44, 0x78)
    bg.line.color.rgb = RGBColor(0x35, 0x5A, 0x88); bg.line.width = Pt(0.5)
    bg.shadow.inherit = False
    add_text(s, left + Inches(0.25), lc_y + Inches(0.25), lc_w - Inches(0.5), Inches(0.3),
             k, size=10, bold=True, color=GW_ORANGE)
    size = 13 if '\n' not in v else 10
    add_text(s, left + Inches(0.25), lc_y + Inches(0.55), lc_w - Inches(0.5), lc_h - Inches(0.6),
             v, size=size, bold=True, color=GW_WHITE, line_spacing=1.3)

# Footer date
add_text(s, 0, Inches(6.8), SLIDE_W, Inches(0.35),
         "18 April 2026   ·   Submitted for Guidewire DEVTrails Phase 3",
         size=11, color=RGBColor(0x8A, 0xA0, 0xC0), align=PP_ALIGN.CENTER)

# =========================================================================
# SAVE
# =========================================================================
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        'gigshield-pitch.pptx')
prs.save(out_path)
print(f"Wrote {out_path}")
print(f"Slides: {len(prs.slides)}")
